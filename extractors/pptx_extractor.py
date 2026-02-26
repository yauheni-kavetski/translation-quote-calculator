import io
from pptx import Presentation

def extract_pptx(file_bytes: bytes) -> str:
    """
    Извлекает текст из слайдов pptx.
    """
    import tempfile, os
    content = ""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    content += '\n'.join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    content += "\n"
    finally:
        os.remove(tmp_path)

    return content